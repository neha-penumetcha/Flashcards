"""
extractor.py
Handles pulling raw text out of ANY uploaded file type:
- PDFs (both normal text PDFs AND scanned/image PDFs via OCR fallback)
- PowerPoint (.pptx)
- Word (.docx)
- Plain text (.txt)
- Images (.png/.jpg/.jpeg) - handwritten/scanned notes via OCR

For "bulk material of any size", this is called once per uploaded file,
and returns raw text that later gets chunked (see chunker.py).
"""

import pdfplumber
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
from pdf2image import convert_from_bytes
import io


def extract_from_pdf(file_bytes: bytes) -> str:
    """
    Try normal text extraction first (fast, free, no OCR needed).
    If a page has NO extractable text (common with scanned notes),
    fall back to OCR for just that page.
    """
    text_parts = []
    ocr_needed_pages = []

    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(page_text)
            else:
                # Mark this page index for OCR fallback
                text_parts.append(None)
                ocr_needed_pages.append(i)

    # Only run OCR (slower) on pages that actually need it
    if ocr_needed_pages:
        images = convert_from_bytes(file_bytes, dpi=200)
        for idx in ocr_needed_pages:
            if idx < len(images):
                ocr_text = pytesseract.image_to_string(images[idx])
                text_parts[idx] = ocr_text

    # Filter out any leftover None values and join
    return "\n".join([t for t in text_parts if t])


def extract_from_docx(file_bytes: bytes) -> str:
    """Extract text from Word documents, including tables."""
    doc = Document(io.BytesIO(file_bytes))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]

    # Also grab text inside tables (often skipped by basic extractors)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)

    return "\n".join(parts)


def extract_from_pptx(file_bytes: bytes) -> str:
    """Extract text from every slide, including speaker notes."""
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        parts.append(line)
        # Speaker notes often contain exam-relevant explanations
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text
            if notes.strip():
                parts.append(notes)
    return "\n".join(parts)


def extract_from_image(file_bytes: bytes) -> str:
    """OCR for handwritten/scanned single images."""
    image = Image.open(io.BytesIO(file_bytes))
    return pytesseract.image_to_string(image)


def extract_from_txt(file_bytes: bytes) -> str:
    return file_bytes.decode("utf-8", errors="ignore")


def extract_text(filename: str, file_bytes: bytes) -> str:
    """
    Single entry point - dispatches to the right extractor based on
    file extension. Call this from app.py for every uploaded file.
    """
    ext = filename.lower().split(".")[-1]

    if ext == "pdf":
        return extract_from_pdf(file_bytes)
    elif ext == "docx":
        return extract_from_docx(file_bytes)
    elif ext == "pptx":
        return extract_from_pptx(file_bytes)
    elif ext in ("png", "jpg", "jpeg"):
        return extract_from_image(file_bytes)
    elif ext == "txt":
        return extract_from_txt(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: .{ext}")
